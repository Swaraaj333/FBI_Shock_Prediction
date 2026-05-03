"""Generate a PPT explaining the model training pipeline."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK = RGBColor(0x0F, 0x0C, 0x29)
BG_MID = RGBColor(0x1A, 0x1A, 0x3E)
ACCENT = RGBColor(0x66, 0x7E, 0xEA)
ACCENT2 = RGBColor(0x76, 0x4B, 0xA2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xA0, 0xA0, 0xCC)
GREEN = RGBColor(0x00, 0xC8, 0x51)
RED = RGBColor(0xFF, 0x44, 0x44)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, left, top, width, height, items,
                     font_size=16, color=LIGHT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)


# ════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)
add_text(slide, 1, 1.5, 11, 1.5,
         "Model Training Pipeline", 44, ACCENT, True, PP_ALIGN.CENTER)
add_text(slide, 1, 3.2, 11, 1,
         "Hemodynamic Shock Predictor — How We Trained the ML Model",
         24, LIGHT, False, PP_ALIGN.CENTER)
add_text(slide, 1, 4.5, 11, 0.8,
         "Built on Real MIMIC-IV ICU Data  •  HistGradientBoostingClassifier  •  86.78% Accuracy",
         18, ACCENT2, False, PP_ALIGN.CENTER)
add_text(slide, 1, 6.2, 11, 0.5,
         "train_model.py", 14, LIGHT, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════
# SLIDE 2 — Training Data Overview
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8,
         "1. Training Data — What We Used", 32, ACCENT, True)
add_bullet_slide(slide, 0.8, 1.5, 5.5, 5, [
    "Source: Real MIMIC-IV ICU Database (PhysioNet)",
    "Raw Files: chartevents.csv, icustays.csv, patients.csv",
    "Processed by: process_mimic_data.py (7-step ETL pipeline)",
    "Output: data/processed_mimic.csv",
    "",
    "Dataset Stats:",
    "  • 6,685 hourly patient records",
    "  • 140 unique ICU stays",
    "  • 72-hour window per patient",
    "  • Vitals sampled every hour",
], 17, LIGHT)
add_bullet_slide(slide, 7, 1.5, 5.5, 5, [
    "Label Distribution:",
    "  • Stable (0): 4,778 rows  (71.5%)",
    "  • Shock  (1): 1,907 rows  (28.5%)",
    "",
    "Label uses 2-hour LOOKAHEAD:",
    "  → We predict if the patient WILL crash",
    "     in the next 1-2 hours",
    "  → This makes it an early-warning system,",
    "     not a reactive alarm",
], 17, LIGHT)

# ════════════════════════════════════════════════
# SLIDE 3 — Features Used
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8,
         "2. Features — What the Model Sees", 32, ACCENT, True)
add_bullet_slide(slide, 0.8, 1.5, 5.5, 5.5, [
    "Core Vitals (5 features):",
    "  • HeartRate  — bpm",
    "  • SysBP     — Systolic Blood Pressure (mmHg)",
    "  • MAP       — Mean Arterial Pressure (mmHg)",
    "  • RespRate  — Breaths per minute",
    "  • SpO2      — Oxygen Saturation (%)",
    "",
    "Derived Feature (1):",
    "  • ShockIndex = HeartRate / SysBP",
    "    (Normal: 0.5-0.7, Shock: > 1.0)",
], 17, LIGHT)
add_bullet_slide(slide, 7, 1.5, 5.5, 5.5, [
    "Temporal Trend Features (4):",
    "  • HR_Change    — Heart Rate change from last hour",
    "  • SysBP_Change — BP change from last hour",
    "  • MAP_Change   — MAP change from last hour",
    "  • SpO2_Change  — O2 change from last hour",
    "",
    "Demographics (2):",
    "  • Age",
    "  • Gender_M  (1 = Male, 0 = Female)",
    "",
    "Total: 12 features per prediction",
], 17, LIGHT)

# ════════════════════════════════════════════════
# SLIDE 4 — Why HistGradientBoosting
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8,
         "3. Model Choice — Why HistGradientBoosting?", 32, ACCENT, True)
add_bullet_slide(slide, 0.8, 1.5, 5.5, 5, [
    "Why NOT Random Forest?",
    "  → Random Forest CANNOT handle NaN values",
    "  → ICU data has frequent missing values",
    "     (sensor disconnections, skipped readings)",
    "  → Would require imputation which can",
    "     hallucinate healthy vitals for sick patients",
    "",
    "Why HistGradientBoosting?",
    "  → Native NaN handling (no imputation needed!)",
    "  → Learns complex feature interactions",
    "  → Handles class imbalance well with weights",
], 17, LIGHT)
add_bullet_slide(slide, 7, 1.5, 5.5, 5, [
    "Hyperparameters:",
    "  • max_iter = 100    (100 boosting rounds)",
    "  • max_depth = 12    (deep trees for complex",
    "                        ICU vital interactions)",
    "  • min_samples_leaf = 5  (prevents overfitting",
    "                            on rare edge cases)",
    "  • random_state = 42     (reproducibility)",
    "",
    "From sklearn.ensemble:",
    "  HistGradientBoostingClassifier",
], 17, LIGHT)

# ════════════════════════════════════════════════
# SLIDE 5 — Preprocessing Tricks
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8,
         "4. Preprocessing — Smart Data Handling", 32, ACCENT, True)
add_bullet_slide(slide, 0.8, 1.5, 5.5, 5, [
    "NaN Masking Strategy:",
    "  • We intentionally mask 15% of Age & Gender",
    "    during training",
    "  • This teaches the model to handle missing",
    "    demographics gracefully at prediction time",
    "  • Vitals are NOT masked — they flow in with",
    "    their natural NaN patterns",
    "",
    "No Global Imputation:",
    "  • We do NOT replace NaN with median/mean",
    "  • The model handles NaN natively",
    "  • This prevents hallucinating healthy vitals",
], 17, LIGHT)
add_bullet_slide(slide, 7, 1.5, 5.5, 5, [
    "Sample Weighting (Critical!):",
    "",
    "  Hypoxia (SpO2 < 90):",
    "    → Only ~56 rows out of 6,685",
    "    → But clinically life-threatening",
    "    → Weight: 50x normal",
    "",
    "  Severe Respiratory Distress:",
    "    (RespRate > 25 AND HeartRate > 110)",
    "    → Weight: 20x normal",
    "",
    "  This forces the model to learn rare but",
    "  critical distress patterns",
], 17, LIGHT)

# ════════════════════════════════════════════════
# SLIDE 6 — Train/Test Split
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8,
         "5. Training Process", 32, ACCENT, True)
add_bullet_slide(slide, 0.8, 1.5, 11, 5, [
    "Step 1: Load processed_mimic.csv  (6,685 rows × 12 features)",
    "",
    "Step 2: Split 80/20  →  Training: 5,348 rows  |  Testing: 1,337 rows",
    "         • Stratified split ensures both sets have same shock/stable ratio",
    "         • random_state=42 for reproducibility",
    "",
    "Step 3: Train HistGradientBoostingClassifier",
    "         • model.fit(X_train, y_train, sample_weight=w_train)",
    "         • Sample weights boost rare critical cases (hypoxia 50x, resp distress 20x)",
    "",
    "Step 4: Evaluate on test set (never seen during training)",
    "",
    "Step 5: Save model artifacts to model/ folder:",
    "         • shock_rf_model.pkl      — the trained model (381 KB)",
    "         • feature_names.pkl       — ordered list of 12 feature names",
], 17, LIGHT)

# ════════════════════════════════════════════════
# SLIDE 7 — Results
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8,
         "6. Results — Model Performance", 32, ACCENT, True)
add_text(slide, 3.5, 1.8, 6, 0.8,
         "Overall Accuracy: 86.78%", 36, GREEN, True, PP_ALIGN.CENTER)
add_bullet_slide(slide, 0.8, 3, 5.5, 4, [
    "Classification Report:",
    "",
    "              Precision  Recall   F1",
    "  Stable:      0.88      0.95    0.91",
    "  Shock:       0.83      0.67    0.74",
    "",
    "Precision = When it says SHOCK, it's right 83%",
    "Recall    = It catches 67% of actual shocks",
], 16, LIGHT)
add_bullet_slide(slide, 7, 3, 5.5, 4, [
    "Sanity Check Results:",
    "",
    "  ✅ Healthy vitals (HR=72, SysBP=120)",
    "     → Correctly predicts STABLE",
    "",
    "  ✅ Shock vitals (HR=140, SysBP=60)",
    "     → Correctly predicts SHOCK",
    "",
    "  ✅ All 7 pytest validation tests PASS",
], 16, LIGHT)

# ════════════════════════════════════════════════
# SLIDE 8 — Shock Labeling Rules
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8,
         "7. How We Define \"Shock\" — Clinical Rules", 32, ACCENT, True)
add_bullet_slide(slide, 0.8, 1.5, 11, 5.5, [
    "A patient-hour is labeled SHOCK (1) if ANY of these conditions are true:",
    "",
    "  Rule 1:  MAP < 65 mmHg              — Hemodynamic instability",
    "  Rule 2:  Shock Index > 1.0           — Tachycardia relative to blood pressure",
    "  Rule 3:  SysBP < 90 mmHg            — Frank hypotension",
    "  Rule 4:  SpO2 < 90%                  — Hypoxia (oxygen deprivation)",
    "  Rule 5:  RespRate > 25 AND HR > 110  — Combined respiratory distress",
    "",
    "The FINAL label uses a 2-hour LOOKAHEAD (shift -2):",
    "  → We don't label what IS happening now",
    "  → We label what WILL happen 2 hours from now",
    "  → This is what makes this an EARLY WARNING system",
    "",
    "These rules come from standard ICU clinical guidelines (Surviving Sepsis Campaign).",
], 17, LIGHT)

# ════════════════════════════════════════════════
# SLIDE 9 — Pipeline Summary
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8,
         "8. End-to-End Pipeline Summary", 32, ACCENT, True)
add_bullet_slide(slide, 0.8, 1.5, 11, 5.5, [
    "Raw MIMIC-IV Data (PhysioNet)",
    "     ↓  process_mimic_data.py",
    "Processed Dataset (6,685 rows × 12 features + labels)",
    "     ↓  train_model.py",
    "Trained Model (shock_rf_model.pkl — 381 KB)",
    "     ↓  flask_api.py",
    "REST API on localhost:5000 (/predict, /predict_csv)",
    "     ↓  streamlit_app.py",
    "Interactive UI on localhost:8501",
    "",
    "Key Files:",
    "  • ml_pipeline/process_mimic_data.py  — ETL (raw → clean)",
    "  • ml_pipeline/train_model.py         — Training (clean → model)",
    "  • app/flask_api.py                   — Inference API",
    "  • app/streamlit_app.py               — Frontend UI",
    "  • tests/test_app.py                  — 7 validation tests (all passing)",
], 17, LIGHT)

# ════════════════════════════════════════════════
# SLIDE 10 — Thank You
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, 1, 2.5, 11, 1.5,
         "Thank You", 48, ACCENT, True, PP_ALIGN.CENTER)
add_text(slide, 1, 4.2, 11, 1,
         "Hemodynamic Shock Predictor — Model Training Pipeline",
         22, LIGHT, False, PP_ALIGN.CENTER)
add_text(slide, 1, 5.5, 11, 0.8,
         "Built with MIMIC-IV  •  HistGradientBoosting  •  86.78% Accuracy  •  12 Features",
         16, ACCENT2, False, PP_ALIGN.CENTER)

# Save
out_path = os.path.join(os.path.dirname(__file__), "Model_Training_Pipeline.pptx")
prs.save(out_path)
print(f"PPT saved to: {out_path}")
